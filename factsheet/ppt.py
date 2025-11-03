from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CATEGORY_TYPE
from typing import Optional
from pathlib import Path


def open_presentation(path: str):
    return Presentation(path)


def _iter_shapes(shapes):
    for shp in shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shp.shapes)


def set_text_by_name(slides, name: str, text: str) -> bool:
    target_name = str(name)
    for slide in slides:
        for shp in _iter_shapes(slide.shapes):
            if getattr(shp, "name", None) == target_name and hasattr(shp, "text_frame"):
                tf = shp.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    # Preserve formatting by editing first run in place
                    p = tf.paragraphs[0]
                    r0 = p.runs[0]
                    r0.text = text
                    # Clear remaining runs' text
                    for r in p.runs[1:]:
                        r.text = ""
                else:
                    # Fallback: create run (format may default)
                    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                    run = p.add_run()
                    run.text = text
                return True
    return False


def replace_picture_by_name(slides, name: str, path: str) -> bool:
    for slide in slides:
        for shp in slide.shapes:
            if shp.name == name:
                left, top, width, height = shp.left, shp.top, shp.width, shp.height
                slide.shapes._spTree.remove(shp._element)
                slide.shapes.add_picture(path, left, top, width=width, height=height)
                return True
    return False


def list_shape_names(prs: Presentation) -> list:
    report = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            report.append({
                "slide": idx,
                "name": getattr(shp, 'name', ''),
                "type": shp.shape_type,
                "has_text": hasattr(shp, 'text_frame')
            })
    return report


def add_picture_slide(prs: Presentation, image_path: str, title_text: str | None = None):
    # Use blank layout (usually index 6); fallback to 0 if out of range
    layout_idx = 6 if len(prs.slide_layouts) > 6 else 0
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Optional title textbox at top
    if title_text:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    # Add image fitted within margins
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(9.0)
    height = None
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    return slide


def add_textbox_slide(prs: Presentation, title: str, body: str):
    layout_idx = 6 if len(prs.slide_layouts) > 6 else 0
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    tx_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
    tf_t = tx_title.text_frame
    tf_t.clear()
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(20)
    p_t.font.bold = True

    tx_body = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9.0), Inches(1.5))
    tf_b = tx_body.text_frame
    tf_b.clear()
    p_b = tf_b.paragraphs[0]
    p_b.text = body
    p_b.font.size = Pt(14)
    return slide


def _find_shape(slides, name: str):
    for slide in slides:
        for shp in slide.shapes:
            if getattr(shp, 'name', None) == name:
                return slide, shp
    return None, None


def _shape_bounds(shp):
    return shp.left, shp.top, shp.width, shp.height


def _remove_shape(slide, shp):
    slide.shapes._spTree.remove(shp._element)


def write_table_by_name(slides, name: str, df) -> bool:
    """Write a pandas.DataFrame into a named shape area.

    - If the named shape is a table and dimensions match, clear and refill.
    - Otherwise, remove the old shape and create a new table with df dimensions
      inside the same bounds.
    """
    slide, shp = _find_shape(slides, name)
    if shp is None:
        return False

    left, top, width, height = _shape_bounds(shp)
    rows = len(df) + 1
    cols = len(df.columns)

    # If shape is a table and dimensions match → clear & refill
    if hasattr(shp, 'has_table') and shp.has_table:
        tbl = shp.table
        # Determine overlapping columns only to preserve formatting
        table_cols = tbl.columns.__len__()
        use_cols = min(table_cols, cols)

        def _set_cell_text_preserve(cell, value: str):
            tf = cell.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                p = tf.paragraphs[0]
                r0 = p.runs[0]
                r0.text = value
            else:
                # Fall back to appending a run (format likely defaults)
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                run = p.add_run()
                run.text = value

        # Header
        for j in range(use_cols):
            _set_cell_text_preserve(tbl.cell(0, j), str(list(df.columns)[j]))

        # Write rows up to table capacity (excluding header)
        table_rows_capacity = tbl.rows.__len__() - 1
        use_rows = min(table_rows_capacity, len(df))
        for i in range(use_rows):
            row_vals = list(df.iloc[i])
            for j in range(use_cols):
                _set_cell_text_preserve(tbl.cell(i + 1, j), str(row_vals[j]))

        # Blank any remaining rows
        for i in range(use_rows, table_rows_capacity):
            for j in range(use_cols):
                _set_cell_text_preserve(tbl.cell(i + 1, j), "")
        return True

    else:
        # Not a table → remove placeholder box
        _remove_shape(slide, shp)

    # Create new table in same bounds
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, col in enumerate(df.columns):
        tbl.cell(0, j).text = str(col)
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        for j, val in enumerate(row):
            tbl.cell(i, j).text = str(val)
    return True


def replace_native_chart_data_by_name(slides, name: str, categories, series_dict, number_format: str | None = None, percent_y_axis: bool = False, force_text_categories: bool = False, date_categories: bool = False, date_format: str | None = None) -> bool:
    """Update a native PowerPoint chart's embedded workbook data in place.

    - Preserves all chart formatting (colors, fonts, gridlines, axes) in the template.
    - categories: list[str]
    - series_dict: mapping series name -> list[float]
    - number_format: optional number format like "0.00%" for values
    """
    slide, shp = _find_shape(slides, name)
    if shp is None or not getattr(shp, 'has_chart', False):
        return False

    chart = shp.chart
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, values in series_dict.items():
        chart_data.add_series(series_name, values)

    try:
        chart.replace_data(chart_data)
    except Exception:
        # Likely an externally-linked chart workbook; signal failure so caller can fallback
        return False

    if number_format is not None:
        for series in chart.series:
            try:
                series.data_labels.number_format = number_format
            except Exception:
                pass
        try:
            # Apply to value axis tick labels to ensure display as percent
            chart.value_axis.tick_labels.number_format = number_format
        except Exception:
            pass

    if percent_y_axis:
        try:
            chart.value_axis.tick_labels.number_format = "0%"
        except Exception:
            pass

    if force_text_categories:
        try:
            chart.category_axis.category_type = XL_CATEGORY_TYPE.CATEGORY
        except Exception:
            pass
    if date_categories:
        try:
            chart.category_axis.category_type = XL_CATEGORY_TYPE.TIME
            if date_format:
                chart.category_axis.tick_labels.number_format = date_format
        except Exception:
            pass
    return True


def get_chart_structure(slides, name: str):
    """Return basic structure of a native chart if found.

    Returns dict with keys: categories (list[str]|None), series_names (list[str]), points_per_series (int|None)
    or None if shape not a native chart.
    """
    slide, shp = _find_shape(slides, name)
    if shp is None or not getattr(shp, 'has_chart', False):
        return None
    chart = shp.chart
    categories = None
    try:
        cats = []
        for c in chart.plots[0].categories:
            try:
                cats.append(str(c.label))
            except Exception:
                cats.append(str(c))
        categories = cats if cats else None
    except Exception:
        categories = None
    series_names = []
    points_per_series = None
    try:
        for s in chart.series:
            series_names.append(s.name)
            if points_per_series is None:
                try:
                    points_per_series = len(list(s.values))
                except Exception:
                    points_per_series = None
    except Exception:
        pass
    return {
        'categories': categories,
        'series_names': series_names,
        'points_per_series': points_per_series,
    }

